"""
Vues pour l'application d'analyse
"""
import json
import io
from decimal import Decimal
from datetime import datetime
from django.shortcuts import render
from django.contrib.auth.decorators import login_required
from django.http import HttpResponseForbidden, HttpResponse
from .models import ChartConfig

# Imports pour les exports
from reportlab.lib.pagesizes import A4, letter
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from pptx import Presentation
from pptx.util import Inches, Pt
import matplotlib
matplotlib.use('Agg')  # Backend sans interface graphique
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches


def convert_decimals(obj):
    """Convertit récursivement les Decimal en float pour la sérialisation JSON"""
    if isinstance(obj, Decimal):
        return float(obj)
    elif isinstance(obj, dict):
        return {key: convert_decimals(value) for key, value in obj.items()}
    elif isinstance(obj, list):
        return [convert_decimals(item) for item in obj]
    return obj


def user_can_access_analysis(user):
    """
    Vérifie si l'utilisateur peut accéder aux analyses.
    Accessible uniquement aux ADMIN, EMPLOYEE et VOLUNTEER_GOVERNANCE.
    """
    if user.is_superuser:
        return True

    try:
        volunteer = user.volunteer_profile
        return volunteer and volunteer.role in ['ADMIN', 'EMPLOYEE', 'VOLUNTEER_GOVERNANCE']
    except AttributeError:
        return False


@login_required
def analysis_dashboard(request):
    """
    Page principale d'analyse avec tous les graphiques configurés.
    Accessible uniquement aux utilisateurs autorisés.
    """
    # Vérification des permissions
    if not user_can_access_analysis(request.user):
        return HttpResponseForbidden(
            "Vous n'avez pas les permissions nécessaires pour accéder aux analyses. "
            "Cette section est réservée aux administrateurs, salariés et bénévoles de gouvernance."
        )

    # Récupérer tous les graphiques actifs groupés par section
    charts = ChartConfig.objects.filter(is_active=True).order_by('section', 'display_order', 'title')

    # Grouper par section
    sections_data = {}
    for chart in charts:
        data = chart.get_chart_data()

        chart_info = {
            'id': chart.id,
            'title': chart.title,
            'description': chart.description,
            'chart_type': chart.chart_type,
            'size': chart.size,
            'y_axis_label': chart.y_axis_label,
            'x_axis_label': chart.x_axis_label,
            'data_json': json.dumps(convert_decimals(data)),
            'has_error': 'error' in data
        }

        if 'error' in data:
            chart_info['error_message'] = data['error']

        # Grouper par section
        section_key = chart.section
        section_label = chart.get_section_display()

        if section_key not in sections_data:
            sections_data[section_key] = {
                'label': section_label,
                'charts': []
            }

        sections_data[section_key]['charts'].append(chart_info)

    # Ordre des sections pour l'affichage
    section_order = ['IMPACT', 'DEMOGRAPHIC', 'FINANCIAL', 'OPERATIrosaL', 'TRENDS', 'ADVANCED']
    ordered_sections = [
        {'key': key, 'label': sections_data[key]['label'], 'charts': sections_data[key]['charts']}
        for key in section_order if key in sections_data
    ]

    context = {
        'sections': ordered_sections,
        'can_export': True,  # Pour activer les boutons d'export plus tard
    }

    return render(request, 'analysis/dashboard.html', context)


def generate_chart_image(chart_data, chart_type, title):
    """Génère une image PNG d'un graphique à partir des données Chart.js"""
    fig, ax = plt.subplots(figsize=(10, 6))

    try:
        datasets = chart_data.get('datasets', [])
        labels = chart_data.get('labels', [])

        if chart_type in ['pie', 'doughnut']:
            # Graphique circulaire
            if datasets and datasets[0].get('data'):
                data = datasets[0]['data']
                colors = datasets[0].get('backgroundColor', [])
                ax.pie(data, labels=labels, autopct='%1.1f%%', colors=colors, startangle=90)
                ax.axis('equal')

        elif chart_type == 'bar':
            # Graphique à barres
            x = range(len(labels))
            width = 0.8 / len(datasets) if len(datasets) > 1 else 0.6

            for i, dataset in enumerate(datasets):
                offset = (i - len(datasets) / 2) * width + width / 2
                bars = ax.bar([p + offset for p in x], dataset['data'], width,
                             label=dataset.get('label', f'Série {i+1}'),
                             color=dataset.get('backgroundColor', None))

            ax.set_xticks(x)
            ax.set_xticklabels(labels, rotation=45, ha='right')
            ax.legend()
            ax.grid(axis='y', alpha=0.3)

        elif chart_type == 'line':
            # Graphique en courbes
            for dataset in datasets:
                ax.plot(labels, dataset['data'],
                       label=dataset.get('label', ''),
                       marker='o',
                       color=dataset.get('borderColor', None))

            ax.legend()
            ax.grid(True, alpha=0.3)
            plt.xticks(rotation=45, ha='right')

        elif chart_type == 'stacked_bar':
            # Graphique à barres empilées
            bottom = [0] * len(labels)
            for dataset in datasets:
                ax.bar(labels, dataset['data'], label=dataset.get('label', ''),
                      bottom=bottom, color=dataset.get('backgroundColor', None))
                bottom = [b + d for b, d in zip(bottom, dataset['data'])]

            ax.legend()
            ax.grid(axis='y', alpha=0.3)
            plt.xticks(rotation=45, ha='right')

        ax.set_title(title, fontsize=14, fontweight='bold', pad=20)
        plt.tight_layout()

        # Sauvegarder en buffer
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', dpi=150, bbox_inches='tight')
        buffer.seek(0)
        plt.close(fig)

        return buffer

    except Exception as e:
        plt.close(fig)
        print(f"Error generating chart image: {e}")
        return None


@login_required
def export_pdf(request):
    """Export PDF des analyses"""
    if not user_can_access_analysis(request.user):
        return HttpResponseForbidden("Vous n'avez pas les permissions nécessaires.")

    # Créer le PDF en mémoire
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=0.5*inch, bottomMargin=0.5*inch)
    story = []
    styles = getSampleStyleSheet()

    # Style personnalisé pour le titre
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor='#1f2937',
        spaceAfter=30,
        alignment=TA_CENTER
    )

    # Page de titre
    story.append(Paragraph("Analyses & Indicateurs", title_style))
    story.append(Paragraph(f"rosa - {datetime.now().strftime('%d/%m/%Y')}", styles['Normal']))
    story.append(Spacer(1, 0.5*inch))

    # Récupérer les graphiques
    charts = ChartConfig.objects.filter(is_active=True).order_by('section', 'display_order')

    current_section = None
    for chart in charts:
        # Nouvelle section
        if chart.section != current_section:
            if current_section is not None:
                story.append(PageBreak())
            story.append(Paragraph(chart.get_section_display(), styles['Heading1']))
            story.append(Spacer(1, 0.2*inch))
            current_section = chart.section

        # Titre du graphique
        story.append(Paragraph(chart.title, styles['Heading2']))
        if chart.description:
            story.append(Paragraph(chart.description, styles['Normal']))
        story.append(Spacer(1, 0.1*inch))

        # Générer et ajouter l'image du graphique
        data = chart.get_chart_data()
        if 'error' not in data:
            img_buffer = generate_chart_image(data, chart.chart_type, chart.title)
            if img_buffer:
                img = Image(img_buffer, width=6*inch, height=3.6*inch)
                story.append(img)
        else:
            story.append(Paragraph(f"Erreur: {data['error']}", styles['Normal']))

        story.append(Spacer(1, 0.3*inch))

    # Construire le PDF
    doc.build(story)
    buffer.seek(0)

    # Retourner la réponse
    response = HttpResponse(buffer.read(), content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename="analyses_rosa_{datetime.now().strftime("%Y%m%d")}.pdf"'
    return response


@login_required
def export_ppt(request):
    """Export PowerPoint des analyses"""
    if not user_can_access_analysis(request.user):
        return HttpResponseForbidden("Vous n'avez pas les permissions nécessaires.")

    # Créer la présentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide de titre
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Analyses & Indicateurs"
    subtitle.text = f"rosa - {datetime.now().strftime('%d/%m/%Y')}"

    # Récupérer les graphiques
    charts = ChartConfig.objects.filter(is_active=True).order_by('section', 'display_order')

    current_section = None
    for chart in charts:
        # Slide de section
        if chart.section != current_section:
            section_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
            txBox = section_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
            tf = txBox.text_frame
            tf.text = chart.get_section_display()
            p = tf.paragraphs[0]
            p.font.size = Pt(44)
            p.font.bold = True
            current_section = chart.section

        # Slide avec graphique
        blank_slide_layout = prs.slide_layouts[5]  # Blank
        slide = prs.slides.add_slide(blank_slide_layout)

        # Titre
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = txBox.text_frame
        tf.text = chart.title
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True

        # Description
        if chart.description:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
            tf = txBox.text_frame
            tf.text = chart.description
            p = tf.paragraphs[0]
            p.font.size = Pt(14)

        # Graphique
        data = chart.get_chart_data()
        if 'error' not in data:
            img_buffer = generate_chart_image(data, chart.chart_type, chart.title)
            if img_buffer:
                pic = slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.5),
                                              width=Inches(8), height=Inches(5))
        else:
            # Afficher l'erreur
            txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
            tf = txBox.text_frame
            tf.text = f"Erreur: {data['error']}"
            p = tf.paragraphs[0]
            p.font.size = Pt(16)

    # Sauvegarder en mémoire
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    # Retourner la réponse
    response = HttpResponse(buffer.read(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = f'attachment; filename="analyses_rosa_{datetime.now().strftime("%Y%m%d")}.pptx"'
    return response
